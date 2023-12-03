#!/usr/bin/env python
import os
import io
import tempfile
from subprocess import call
from pdf2image import convert_from_path
from pptx import Presentation
import openai
import requests
from PIL import Image
import requests
import base64


## Sometimes ffmpeg is avconv
FFMPEG_NAME = 'ffmpeg'
#FFMPEG_NAME = 'avconv'

pptx_path='example/testlong.pptx'
output_path='example/test.mp4'

# Function to generate and save audio
def generate_and_save_audio(text, audio_path):
     # Retrieve API key from environment variable
    openai_api_key = os.getenv('OPENAI_API_KEY')

    if openai_api_key is None:
        raise ValueError("Please set the OPENAI_API_KEY environment variable.")
    
    headers = {
        'Authorization': f'Bearer {openai_api_key}',
        'Content-Type': 'application/json',
    }
    
    data = {
        "model": "tts-1-hd",
        "input": text,
        "voice": "echo",
        "speed": 1.1,
        "response_format": "aac"
    }

    # POST request to OpenAI's API
    response = requests.post('https://api.openai.com/v1/audio/speech', headers=headers, json=data)
    
    if response.status_code != 200:
        raise Exception(f"Request failed with status code {response.status_code}: {response.text}")
    
    # Save the binary content to the file
    with open(audio_path, 'wb') as f:
        f.write(response.content)  # Note it's 'content' not 'data' when dealing with requests library

    return audio_path

# Function to encode image to base64, handling PNG with transparency
def encode_image(image):
    buffered = io.BytesIO()
    if image.mode in ("RGBA", "LA"):
        background = Image.new(image.mode[:-1], image.size, (255, 255, 255))
        background.paste(image, image.split()[-1])
        image = background.convert("RGB")
    image.save(buffered, format="JPEG")
    return base64.b64encode(buffered.getvalue()).decode('utf-8')

# Function to make API call to OpenAI
def ask_openai(base64_image, notes, prev_slide_transcript=None):
     # Retrieve API key from environment variable
    openai_api_key = os.getenv('OPENAI_API_KEY')

    if openai_api_key is None:
        raise ValueError("Please set the OPENAI_API_KEY environment variable.")
    
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {openai_api_key}"
    }
    payload = {
        "model": "gpt-4-vision-preview",
        "messages": [
            {
                "role": "system",
                "content": "You are an engaging and knowledgable presenter of a slide presentation. Given the slide content and speaker notes, generate a spoken narration for the slide. You transcript should be short, no more than 2-3 paragraphs, and in the style and quality of a good TED talk. The transcript should be only raw text, do not include markdown or URLs. It should only have what a normal person would speak aloud, in easy to understand language."
            },
        ],
        "max_tokens": 500
    }

    if prev_slide_transcript is not None:
        payload['messages'].append(
        {
            "role": "user",
            "content": [
                {"type": "text", "text": f"Previous slides transcript: \n---\n{prev_slide_transcript}\n---"}
            ]
        })
    else:
         payload['messages'].append(
        {
            "role": "user",
            "content": [
                {"type": "text", "text": f"This is the beginning of the presentation."}
            ]
        })
    payload['messages'].append(
        {
            "role": "user",
            "content": [
                {"type": "text", "text": f"Use this presentation slide along with the speaker notes to create your voiceover. Do not include a welcome or greeting unless it is the beginning of the presentation. Speaker notes: {notes}"},
                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
            ]
        })
    response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload)
    json = response.json()
    if 'choices' in json:
        return json['choices'][0]['message']['content']
    else:
        raise Exception(f"Request failed with status code {response.status_code}: {response.text}")

def main():
    full_script = ''
    with tempfile.TemporaryDirectory() as temp_path:
        call(['soffice', '--headless', '--convert-to', 'pdf:writer_pdf_Export', pptx_path, '--outdir', temp_path])
        pdf_file_name = os.path.splitext(os.path.basename(pptx_path))[0] + '.pdf'
        pdf_path = os.path.join(temp_path, pdf_file_name)
        images_from_path = convert_from_path(pdf_path)
        prs = Presentation(pptx_path)
        assert len(images_from_path) == len(prs.slides)
        for i, (slide, image) in enumerate(zip(prs.slides, images_from_path)):
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
            else:
                notes = ''
            image_path = os.path.join(temp_path, 'frame_{}.jpg'.format(i))
            audio_path = os.path.join(temp_path, 'frame_{}.aac'.format(i))

            transcript = ask_openai(encode_image(image), notes, full_script)
            full_script += transcript
            print('SLIDE',i,transcript)
            image.save(image_path)

            generate_and_save_audio(transcript, audio_path)

            ffmpeg_call(image_path, audio_path, temp_path, i)

        video_list = [os.path.join(temp_path, 'frame_{}.ts'.format(i)) \
                        for i in range(len(images_from_path))]
        video_list_str = 'concat:' + '|'.join(video_list)
        ffmpeg_concat(video_list_str, output_path)
        print(full_script)


def ffmpeg_call(image_path, audio_path, temp_path, i):
    out_path_mp4 = os.path.join(temp_path, 'frame_{}.mp4'.format(i))
    out_path_ts = os.path.join(temp_path, 'frame_{}.ts'.format(i))
    call([FFMPEG_NAME, '-loop', '1', '-y', '-i', image_path, '-i', audio_path, '-shortest', '-fflags', '+shortest', '-max_interleave_delta', '200M',
      '-c:v', 'libx264', '-tune', 'stillimage', '-c:a', 'aac',
      '-b:a', '192k', '-vf', 'scale=-1:1080', out_path_mp4])
    call([FFMPEG_NAME, '-y', '-i', out_path_mp4, '-c', 'copy',
          '-bsf:v', 'h264_mp4toannexb', '-f', 'mpegts', out_path_ts])


def ffmpeg_concat(video_list_str, out_path):
    call([FFMPEG_NAME, '-y', '-f', 'mpegts', '-i', '{}'.format(video_list_str),
          '-c', 'copy', '-bsf:a', 'aac_adtstoasc', out_path])

if __name__ == '__main__':
    main()